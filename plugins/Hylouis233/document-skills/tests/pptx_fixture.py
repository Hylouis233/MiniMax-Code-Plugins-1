# Minimal runnable fixtures for the snippets called out in review: one file per format.
# Each script is self-contained, writes only scratch files into the current directory,
# and exits non-zero on failed assertions. Run from any scratch directory:
#   python pdf_fixture.py   (deps: reportlab, pypdf, pymupdf)
#   python pptx_fixture.py  (deps: python-pptx)
#   python xlsx_fixture.py  (deps: openpyxl)
#   python docx_fixture.py  (deps: python-docx, pymupdf, soffice on PATH)
import sys
import zipfile

from lxml import etree
from pptx import Presentation
from pptx.chart.data import BubbleChartData, ChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches, Pt

failures = []


def check(name, cond, extra=""):
    print(("PASS " if cond else "FAIL ") + name + ((" :: " + str(extra)) if not cond and extra else ""))
    if not cond:
        failures.append(name)


# ---- build the deck ------------------------------------------------------------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
tf = box.text_frame
p = tf.paragraphs[0]
r1 = p.add_run()
r1.text = "old wording"
r2 = p.add_run()
r2.text = " linked part"
r2.font.italic = True
r2.font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
r2.hyperlink.address = "https://example.com/docs"

table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.5), Inches(6), Inches(1))
table = table_shape.table
cell = table.cell(0, 1)
ctf = cell.text_frame
cp = ctf.paragraphs[0]
cr = cp.add_run()
cr.text = "old cell text"
cr.font.bold = True
cr.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)

chart_data = ChartData()
chart_data.categories = ["EU", "US"]
chart_data.add_series("Units", (120, 80))
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(7.2), Inches(2.5), Inches(2), Inches(2),
    chart_data,
).chart
chart.has_title = True
chart.chart_title.text_frame.text = "Units by region"
slide.notes_slide.notes_text_frame.text = "Speaker note: explain the regional split."

prs.save("input.pptx")

# ---- edit.md snippet: single-shape run replace keeps styling and hyperlink -----
prs = Presentation("input.pptx")
old, new = "old wording", "new wording"
candidates = []
for i, s in enumerate(prs.slides):
    for shape in s.shapes:
        if shape.has_text_frame and old in shape.text_frame.text:
            candidates.append((i, shape.name, shape))
assert len(candidates) == 1
_, _, target_shape = candidates[0]
tf = target_shape.text_frame
run_hits = [run for par in tf.paragraphs for run in par.runs if old in run.text]
assert len(run_hits) == 1
run_hits[0].text = run_hits[0].text.replace(old, new, 1)

edited_link = [run for par in tf.paragraphs for run in par.runs if run.hyperlink.address]
check("run replace keeps the other run's hyperlink", len(edited_link) == 1 and edited_link[0].hyperlink.address == "https://example.com/docs")
styled = [run for par in tf.paragraphs for run in par.runs if run.font.italic]
check("run replace keeps sibling run styling", len(styled) == 1 and styled[0].font.color.rgb == RGBColor(0xC0, 0x00, 0x00))
prs.save("edited.pptx")

# ---- edit.md snippet: table cell edited at run level ---------------------------
prs2 = Presentation("input.pptx")
old_cell, new_cell = "old cell text", "new cell text"
tbl = next(sh for sh in prs2.slides[0].shapes if sh.has_table).table
cell = tbl.cell(0, 1)
hits = [run for par in cell.text_frame.paragraphs for run in par.runs if old_cell in run.text]
assert len(hits) == 1, "target is duplicated or split across runs in this cell"
hits[0].text = hits[0].text.replace(old_cell, new_cell, 1)

prs2.save("cell-edited.pptx")
prs3 = Presentation("cell-edited.pptx")
cell3 = next(sh for sh in prs3.slides[0].shapes if sh.has_table).table.cell(0, 1)
after_runs = [run for par in cell3.text_frame.paragraphs for run in par.runs]
check("cell run edit keeps bold", any(r.font.bold for r in after_runs))
check("cell run edit keeps color", any(r.font.color and r.font.color.rgb == RGBColor(0x00, 0x70, 0xC0) for r in after_runs))
check("cell run edit changed the text", cell3.text_frame.text == "new cell text")

# the dangerous variant for contrast: assigning cell.text drops run properties
prs4 = Presentation("input.pptx")
tbl4 = next(sh for sh in prs4.slides[0].shapes if sh.has_table).table
tbl4.cell(0, 1).text = "new cell text"
prs4.save("cell-flattened.pptx")
prs5 = Presentation("cell-flattened.pptx")
flat_runs = [run for par in next(sh for sh in prs5.slides[0].shapes if sh.has_table).table.cell(0, 1).text_frame.paragraphs for run in par.runs]
check("cell.text assignment is proven lossy (negative control)", not any(r.font.bold for r in flat_runs))

# ---- analyze.md snippet: grouped-shape walker ----------------------------------
from pptx.oxml import parse_xml

prs6 = Presentation()
slide6 = prs6.slides.add_slide(prs6.slide_layouts[5])
pic_holder = slide6.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
pic_holder.text_frame.text = "nested member"
pic_holder.text_frame.paragraphs[0].runs[0].font.name = "Grouped Face"

GRP = (
    '<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    "<p:nvGrpSpPr><p:cNvPr id=\"901\" name=\"demo group\"/>"
    "<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>"
    "<p:grpSpPr><a:xfrm><a:off x=\"914400\" y=\"914400\"/>"
    '<a:ext cx="4572000" cy="914400"/>'
    '<a:chOff x="0" y="0"/><a:chExt cx="4572000" cy="914400"/></a:xfrm></p:grpSpPr>'
    "</p:grpSp>"
)


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def cached_numeric_values(series, element_name):
    return [
        node.text for node in
        series._element.xpath(f"./c:{element_name}//c:pt/c:v")
    ]


def extract_slide_content(slide):
    shapes = list(iter_shapes(slide.shapes))
    text = [sh.text_frame.text for sh in shapes if sh.has_text_frame and sh.text_frame.text]
    tables = [
        [[cell.text for cell in row.cells] for row in sh.table.rows]
        for sh in shapes if sh.has_table
    ]
    charts = []
    for sh in shapes:
        if not sh.has_chart:
            continue
        chart = sh.chart
        chart_title = (
            chart.chart_title.text_frame.text
            if chart.has_title and chart.chart_title.has_text_frame else ""
        )
        plots = []
        for plot in chart.plots:
            plot_kind = type(plot).__name__
            if plot_kind in {"XyPlot", "BubblePlot"}:
                series = []
                for item in plot.series:
                    values = {
                        "name": item.name,
                        "x_values": cached_numeric_values(item, "xVal"),
                        "y_values": cached_numeric_values(item, "yVal"),
                    }
                    if plot_kind == "BubblePlot":
                        values["bubble_sizes"] = cached_numeric_values(item, "bubbleSize")
                    series.append(values)
                plots.append({"kind": plot_kind, "series": series})
            else:
                categories = [
                    [str(level) for level in label]
                    for label in plot.categories.flattened_labels
                ]
                series = [
                    {"name": item.name, "values": list(item.values)}
                    for item in plot.series
                ]
                plots.append({"kind": plot_kind, "categories": categories, "series": series})
        charts.append({"title": chart_title, "plots": plots})
    notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
    return {"text": text, "tables": tables, "charts": charts, "notes": notes}


content = extract_slide_content(Presentation("input.pptx").slides[0])
check("content inventory emits body text", any("old wording" in value for value in content["text"]), content)
check("content inventory emits table cell text", content["tables"][0][0][1] == "old cell text", content["tables"])
check(
    "content inventory emits chart title, categories, series, and values",
    content["charts"][0]["title"] == "Units by region"
    and content["charts"][0]["plots"][0]["categories"] == [["EU"], ["US"]]
    and content["charts"][0]["plots"][0]["series"]
    == [{"name": "Units", "values": [120.0, 80.0]}],
    content["charts"],
)
check("content inventory emits notes text", "regional split" in content["notes"], content["notes"])

# XY scatter and bubble plots do not have category/value-series semantics.
xy_prs = Presentation()
xy_slide = xy_prs.slides.add_slide(xy_prs.slide_layouts[6])
xy_data = XyChartData()
xy_series = xy_data.add_series("XY series")
xy_series.add_data_point(1, 2)
xy_series.add_data_point(3, 4)
xy_slide.shapes.add_chart(
    XL_CHART_TYPE.XY_SCATTER,
    Inches(0.5), Inches(0.5), Inches(4), Inches(2.5), xy_data,
)
bubble_data = BubbleChartData()
bubble_series = bubble_data.add_series("Bubble series")
bubble_series.add_data_point(5, 6, 7)
xy_slide.shapes.add_chart(
    XL_CHART_TYPE.BUBBLE,
    Inches(0.5), Inches(3.5), Inches(4), Inches(2.5), bubble_data,
)
xy_prs.save("xy-bubble.pptx")
xy_content = extract_slide_content(Presentation("xy-bubble.pptx").slides[0])
xy_plots = [plot for chart in xy_content["charts"] for plot in chart["plots"]]
check(
    "scatter inventory emits x and y caches without category access",
    any(
        plot["kind"] == "XyPlot"
        and plot["series"][0]["x_values"] == ["1", "3"]
        and plot["series"][0]["y_values"] == ["2", "4"]
        for plot in xy_plots
    ),
    xy_plots,
)
check(
    "bubble inventory emits x, y, and bubble-size caches",
    any(
        plot["kind"] == "BubblePlot"
        and plot["series"][0]["x_values"] == ["5"]
        and plot["series"][0]["y_values"] == ["6"]
        and plot["series"][0]["bubble_sizes"] == ["7"]
        for plot in xy_plots
    ),
    xy_plots,
)


sp_element = slide6.shapes[-1]._element  # the textbox; layout 5 still carries a Title placeholder
group_element = parse_xml(GRP)
sp_element.getparent().replace(sp_element, group_element)
group_element.append(sp_element)

flat = list(iter_shapes(slide6.shapes))
check(
    "walker finds the shape nested in the group",
    any(getattr(sh, "text_frame", None) is not None and sh.text_frame.text == "nested member" for sh in flat),
)
check(
    "top-level shapes list hides the nested member (negative control)",
    not any(getattr(sh, "text_frame", None) is not None and sh.text_frame.text == "nested member" for sh in slide6.shapes),
)
grouped_faces = [
    run.font.name
    for shape in iter_shapes(slide6.shapes)
    if shape.has_text_frame
    for paragraph in shape.text_frame.paragraphs
    for run in paragraph.runs
    if run.font.name
]
check("font triage reaches runs nested in groups", "Grouped Face" in grouped_faces, grouped_faces)

# ---- edit.md locator: candidate collection must recurse into groups ------------
old_w, new_w = "nested member", "renamed member"


def iter_shapes_with_path(shapes, path=""):
    for shape in shapes:
        here = f"{path}/{shape.name}" if path else shape.name
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_with_path(shape.shapes, here)
        else:
            yield here, shape


def iter_text_targets(path, shape):
    if shape.has_text_frame:
        yield path, shape.text_frame
    if shape.has_table:
        for row_index, row in enumerate(shape.table.rows):
            for column_index, cell in enumerate(row.cells):
                yield f"{path}/table[{row_index},{column_index}]", cell.text_frame


candidates = [
    (i, location, text_frame)
    for i, s in enumerate(prs6.slides)
    for p, sh in iter_shapes_with_path(s.shapes)
    for location, text_frame in iter_text_targets(p, sh)
    if old_w in text_frame.text
]
check("locator reaches text inside the group", len(candidates) == 1, [(i, p) for i, p, _ in candidates])
check("locator reports a stable nested path", "/" in candidates[0][1], candidates[0][1])
_, _, target = candidates[0]
target.paragraphs[0].runs[0].text = new_w
prs6.save("group-edited.pptx")
prs_g = Presentation("group-edited.pptx")
found = [sh for sh in iter_shapes(prs_g.slides[0].shapes)
         if getattr(sh, "text_frame", None) is not None and sh.text_frame.text == new_w]
check("group member edit persists after save", len(found) == 1)

table_candidates = [
    (i, location, text_frame)
    for i, s in enumerate(Presentation("input.pptx").slides)
    for p, sh in iter_shapes_with_path(s.shapes)
    for location, text_frame in iter_text_targets(p, sh)
    if old_cell in text_frame.text
]
check("locator reaches wording stored only in a table cell", len(table_candidates) == 1)
check("table-cell locator retains row and column",
      table_candidates[0][1].endswith("/table[0,1]"), table_candidates[0][1])

# ---- analyze.md snippet: per-master, script-aware theme font resolution --------
import xml.etree.ElementTree as ET

prs7 = Presentation("input.pptx")
theme_cache = {}
DRAWINGML = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def read_theme_role(root, role):
    node = root.find(f".//a:{role}Font", DRAWINGML)
    if node is None:
        return {"latin": "", "eastAsia": "", "complexScript": "", "scripts": {}}

    def typeface(tag):
        child = node.find(f"a:{tag}", DRAWINGML)
        return "" if child is None else child.get("typeface", "")

    return {
        "latin": typeface("latin"),
        "eastAsia": typeface("ea"),
        "complexScript": typeface("cs"),
        "scripts": {
            child.get("script"): child.get("typeface", "")
            for child in node.findall("a:font", DRAWINGML) if child.get("script")
        },
    }


def theme_faces_for_slide(slide):
    master_part = slide.slide_layout.slide_master.part
    cache_key = str(master_part.partname)
    if cache_key not in theme_cache:
        theme_part = master_part.part_related_by(RT.THEME)
        root = ET.fromstring(theme_part.blob)
        theme_cache[cache_key] = {
            "major": read_theme_role(root, "major"),
            "minor": read_theme_role(root, "minor"),
        }
    return cache_key, theme_cache[cache_key]


def script_tags(text):
    tags = []
    for character in text:
        codepoint = ord(character)
        if 0x3040 <= codepoint <= 0x30FF:
            tags.append("Jpan")
        elif 0xAC00 <= codepoint <= 0xD7AF:
            tags.append("Hang")
        elif 0x2E80 <= codepoint <= 0x9FFF:
            tags.extend(("Hans", "Hant", "Jpan", "Hang"))
        elif 0x0400 <= codepoint <= 0x052F:
            tags.append("Cyrl")
        elif 0x0590 <= codepoint <= 0x05FF:
            tags.append("Hebr")
        elif 0x0600 <= codepoint <= 0x06FF:
            tags.append("Arab")
        elif 0x0900 <= codepoint <= 0x097F:
            tags.append("Deva")
    return list(dict.fromkeys(tags))


def theme_candidates(role_fonts, text):
    faces = [role_fonts["latin"]]
    tags = script_tags(text)
    if any(tag in ("Hans", "Hant", "Jpan", "Hang") for tag in tags):
        faces.append(role_fonts["eastAsia"])
    if any(tag in ("Arab", "Hebr", "Deva") for tag in tags):
        faces.append(role_fonts["complexScript"])
    faces.extend(role_fonts["scripts"].get(tag, "") for tag in tags)
    return [face for face in dict.fromkeys(faces) if face]


master_name, theme_fonts = theme_faces_for_slide(prs7.slides[0])
check(
    "theme major/minor fonts resolve through the slide master relationship",
    theme_fonts["major"]["latin"] and theme_fonts["minor"]["latin"],
    (master_name, theme_fonts),
)
print("theme fonts:", ascii(theme_fonts))


class StubThemePart:
    def __init__(self, major, minor, east_asian="", complex_script="", scripts=None):
        script_nodes = "".join(
            f'<a:font script="{script}" typeface="{face}"/>'
            for script, face in (scripts or {}).items()
        )
        self.blob = (
            f'<a:theme xmlns:a="{DRAWINGML["a"]}"><a:themeElements><a:fontScheme>'
            f'<a:majorFont><a:latin typeface="{major}"/><a:ea typeface="{east_asian}"/>'
            f'<a:cs typeface="{complex_script}"/>{script_nodes}</a:majorFont>'
            f'<a:minorFont><a:latin typeface="{minor}"/><a:ea typeface="{east_asian}"/>'
            f'<a:cs typeface="{complex_script}"/>{script_nodes}</a:minorFont>'
            f'</a:fontScheme></a:themeElements></a:theme>'
        ).encode()


class StubMasterPart:
    def __init__(self, name, major, minor, **theme_options):
        self.partname = name
        self.theme_part = StubThemePart(major, minor, **theme_options)

    def part_related_by(self, relationship_type):
        assert relationship_type == RT.THEME
        return self.theme_part


def stub_slide(name, major, minor, **theme_options):
    master = type("Master", (), {"part": StubMasterPart(name, major, minor, **theme_options)})()
    layout = type("Layout", (), {"slide_master": master})()
    return type("Slide", (), {"slide_layout": layout})()


_, first_fonts = theme_faces_for_slide(stub_slide("/ppt/slideMasters/one.xml", "Head One", "Body One"))
_, second_fonts = theme_faces_for_slide(stub_slide("/ppt/slideMasters/two.xml", "Head Two", "Body Two"))
check(
    "different slide masters resolve their own theme faces",
    first_fonts["major"]["latin"] == "Head One"
    and second_fonts["major"]["latin"] == "Head Two"
    and first_fonts != second_fonts,
    (first_fonts, second_fonts),
)

_, script_fonts = theme_faces_for_slide(stub_slide(
    "/ppt/slideMasters/scripts.xml", "Latin Theme", "Latin Body",
    east_asian="East Asian Theme", complex_script="Complex Script Theme",
    scripts={"Hans": "Simplified Chinese Theme", "Cyrl": "Cyrillic Theme"},
))
check(
    "CJK inherited-font triage includes east-Asian and script-specific faces",
    {"East Asian Theme", "Simplified Chinese Theme"}
    <= set(theme_candidates(script_fonts["major"], "汉字")),
    theme_candidates(script_fonts["major"], "汉字"),
)
check(
    "non-Latin inherited-font triage uses script mappings",
    "Cyrillic Theme" in theme_candidates(script_fonts["minor"], "текст"),
    theme_candidates(script_fonts["minor"], "текст"),
)

check(
    "unstyled runs report as inherited, not as a concrete face",
    all(run.font.name is None for sh in prs7.slides[0].shapes if sh.has_text_frame for par in sh.text_frame.paragraphs for run in par.runs),
)

font_box = prs7.slides[0].shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
font_paragraph = font_box.text_frame.paragraphs[0]
font_paragraph.font.name = "Paragraph Face"
paragraph_run = font_paragraph.add_run()
paragraph_run.text = "paragraph default"
explicit_run = font_paragraph.add_run()
explicit_run.text = "run override"
explicit_run.font.name = "Run Face"
detected_faces = []
for run in font_paragraph.runs:
    if run.font.name:
        detected_faces.append((run.font.name, "run"))
    elif font_paragraph.font.name:
        detected_faces.append((font_paragraph.font.name, "paragraph defaults"))
check("font triage reports the run face and source", ("Run Face", "run") in detected_faces, detected_faces)
check(
    "font triage reports the paragraph face and source",
    ("Paragraph Face", "paragraph defaults") in detected_faces,
    detected_faces,
)

# ---- analyze.md bounded ZIP/XML health check ----------------------------------
MAX_XML_PART = 20 * 1024 * 1024
MAX_ENTRY = 100 * 1024 * 1024
MAX_TOTAL_UNCOMPRESSED = 500 * 1024 * 1024
MAX_COMPRESSION_RATIO = 200
safe_xml_parser = etree.XMLParser(
    load_dtd=False,
    resolve_entities=False,
    no_network=True,
    huge_tree=False,
    recover=False,
)


def validate_pptx_package(path):
    with zipfile.ZipFile(path) as archive:
        infos = archive.infolist()
        names = {info.filename for info in infos}
        assert len(names) == len(infos)
        assert "[Content_Types].xml" in names and "ppt/presentation.xml" in names
        assert sum(info.file_size for info in infos) <= MAX_TOTAL_UNCOMPRESSED
        actual_total = 0
        for info in infos:
            assert info.file_size <= MAX_ENTRY
            assert info.file_size / max(info.compress_size, 1) <= MAX_COMPRESSION_RATIO
            is_xml = info.filename.endswith((".xml", ".rels"))
            if is_xml:
                assert info.file_size <= MAX_XML_PART
            chunks = []
            actual_size = 0
            with archive.open(info) as stream:
                while chunk := stream.read(64 * 1024):
                    actual_size += len(chunk)
                    actual_total += len(chunk)
                    assert actual_size <= MAX_ENTRY
                    assert actual_total <= MAX_TOTAL_UNCOMPRESSED
                    if is_xml:
                        chunks.append(chunk)
            assert actual_size == info.file_size
            if is_xml:
                etree.fromstring(b"".join(chunks), parser=safe_xml_parser)


try:
    validate_pptx_package("input.pptx")
    healthy_package_passed = True
except Exception:
    healthy_package_passed = False
check("bounded package health check accepts an ordinary PPTX", healthy_package_passed)

with zipfile.ZipFile("compressed-bomb.pptx", "w", zipfile.ZIP_DEFLATED) as archive:
    archive.writestr("[Content_Types].xml", "<Types/>")
    archive.writestr("ppt/presentation.xml", "<presentation>" + (" " * 2_000_000) + "</presentation>")
try:
    validate_pptx_package("compressed-bomb.pptx")
    archive_bomb_rejected = False
except AssertionError:
    archive_bomb_rejected = True
check("PPTX compression bomb is rejected before XML expansion", archive_bomb_rejected)

print("\n" + ("ALL PPTX FIXTURES PASSED" if not failures else f"{len(failures)} FAILURES: {failures}"))
sys.exit(0 if not failures else 1)
